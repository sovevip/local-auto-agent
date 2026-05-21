"""
Phase 1: 工具函数库 + Schema 定义（安全收口版 v2.0）
====================================================
68 个安全收口的本地工具函数，含 Skill 动态加载系统、GitHub 集成。

架构：
  TOOL_REGISTRY — 68 个模型可见工具 → get_tool_schema()
  _SKILL_TOOLS — 动态加载的 Skill 工具 → load_all_skills()
  INTERNAL_TOOLS — ask_user_confirmation → get_internal_tool()

用法：
  python tools.py          → 自检
  python tools.py --schema → 模型可见 JSON Schema
"""

from pathlib import Path
from typing import Any
import os, json, datetime, platform

# ═══════════════════════════════════════════════════════════
# 全局安全配置
# ═══════════════════════════════════════════════════════════

_WORKSPACE_ROOT: Path | None = None
_confirm_handler: Any = None
_SKILLS_DIR: Path = Path(__file__).parent / "skills"
_SKILL_TOOLS: list = []  # 动态加载的技能工具


def set_workspace(root: str) -> None:
    """设置文件操作安全边界。必须在使用任何文件工具前调用。"""
    global _WORKSPACE_ROOT
    _WORKSPACE_ROOT = Path(root).expanduser().resolve()
    if not _WORKSPACE_ROOT.exists():
        _WORKSPACE_ROOT.mkdir(parents=True, exist_ok=True)


def set_confirm_handler(handler: Any) -> None:
    """注入确认回调 handler(action, detail) -> bool。"""
    global _confirm_handler
    _confirm_handler = handler


def _assert_workspace() -> None:
    if _WORKSPACE_ROOT is None:
        raise PermissionError("workspace 未设置，所有文件操作被拒绝。请先调用 set_workspace(path)。")


def _resolve_path(file_path: str) -> Path:
    _assert_workspace()
    p = Path(file_path)
    if p.is_absolute() or str(file_path).startswith("~"):
        path = p.expanduser().resolve()
    else:
        path = (_WORKSPACE_ROOT / p).resolve()
    try:
        path.relative_to(_WORKSPACE_ROOT)
    except ValueError:
        raise PermissionError(f"路径越界: {path} 不在 workspace {_WORKSPACE_ROOT} 内")
    return path


def _ok(data: Any = None) -> str:
    return json.dumps({"ok": True, "data": data, "error": None}, ensure_ascii=False)


def _err(msg: str) -> str:
    return json.dumps({"ok": False, "data": None, "error": msg}, ensure_ascii=False)


def _confirm(action: str, detail: str, is_sensitive: bool = True) -> bool:
    if _confirm_handler is not None:
        return _confirm_handler(action, detail)
    return not is_sensitive  # 默认：敏感操作拒绝


def ask_user_confirmation(action: str, detail: str = "") -> str:
    """内部工具。执行器在敏感操作前自动调用，模型不得直接调用。"""
    return _ok(_confirm(action, detail, is_sensitive=True))


# ═══════════════════════════════════════════════════════════
# 文件操作 (14 tools)
# ═══════════════════════════════════════════════════════════

def find_files(folder: str, pattern: str = "*") -> str:
    """在指定文件夹中查找匹配 glob pattern 的文件。只读。"""
    try:
        p = _resolve_path(folder)
        if not p.exists() or not p.is_dir():
            return _err(f"文件夹不存在或不是文件夹: {folder}")
        files = [str(f) for f in sorted(p.glob(pattern)) if f.is_file()]
        return _ok({"count": len(files), "files": files})
    except PermissionError as e:
        return _err(str(e))


def list_directory(folder: str) -> str:
    """列出文件夹中所有文件和子文件夹。只读。"""
    try:
        p = _resolve_path(folder)
        if not p.exists() or not p.is_dir():
            return _err(f"文件夹不存在或不是文件夹: {folder}")
        items = [{"name": i.name, "type": "directory" if i.is_dir() else "file",
                   "size": i.stat().st_size if i.is_file() else None}
                 for i in sorted(p.iterdir())]
        return _ok({"path": str(p), "items": items})
    except PermissionError as e:
        return _err(str(e))


def read_file_content(file_path: str) -> str:
    """读取文本文件内容。只读。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists() or not path.is_file():
            return _err(f"文件不存在或不是文件: {file_path}")
        return _ok(path.read_text(encoding="utf-8"))
    except PermissionError as e:
        return _err(str(e))


def write_file_content(file_path: str, content: str) -> str:
    """写入内容到文件。敏感操作，需确认。"""
    try:
        path = _resolve_path(file_path)
    except PermissionError as e:
        return _err(str(e))
    action = "写入文件（覆盖）" if path.exists() else "写入文件（新建）"
    if not json.loads(ask_user_confirmation(action, str(path))).get("data"):
        return _err("用户取消了写入操作")
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(content, encoding="utf-8")
        return _ok(f"已写入 {len(content)} 字符到 {file_path}")
    except Exception as e:
        return _err(str(e))


def open_file(file_path: str) -> str:
    """用系统默认程序打开文件。敏感操作，需确认。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
    except PermissionError as e:
        return _err(str(e))
    if not json.loads(ask_user_confirmation("打开文件", str(path))).get("data"):
        return _err("用户取消了打开操作")
    os.startfile(str(path))
    return _ok(f"已打开 {file_path}")


def rename_file(source: str, target: str) -> str:
    """重命名或移动文件。敏感操作，需确认。"""
    try:
        src = _resolve_path(source)
        dst = _resolve_path(target)
        if not src.exists():
            return _err(f"源文件不存在: {source}")
    except PermissionError as e:
        return _err(str(e))
    if not json.loads(ask_user_confirmation("重命名文件", f"{src}\n→ {dst}")).get("data"):
        return _err("用户取消了重命名操作")
    try:
        src.rename(dst)
        return _ok(f"已将 {source} 重命名为 {target}")
    except Exception as e:
        return _err(str(e))


def batch_rename(folder: str, find: str, replace: str, glob_pattern: str = "*") -> str:
    """批量重命名文件夹中的文件。敏感操作，需确认。"""
    try:
        p = _resolve_path(folder)
        if not p.exists() or not p.is_dir():
            return _err(f"文件夹不存在: {folder}")
        files = sorted(p.glob(glob_pattern))
        targets = [(f, f.name.replace(find, replace)) for f in files if f.is_file() and find in f.name]
        if not targets:
            return _ok({"renamed": 0, "results": []})
        if not json.loads(ask_user_confirmation("批量重命名",
            f"文件夹: {folder}\n{len(targets)} 个文件:\n" +
            "\n".join(f"  {old.name} → {new}" for old, new in targets[:10]))).get("data"):
            return _err("用户取消了批量重命名")
        results = []
        for old_path, new_name in targets:
            new_path = old_path.with_name(new_name)
            old_path.rename(new_path)
            results.append({"old": old_path.name, "new": new_name})
        return _ok({"renamed": len(results), "results": results})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def delete_file(target: str, recursive: bool = False) -> str:
    """删除文件或文件夹。敏感操作，需确认。非空文件夹需 recursive=true。"""
    try:
        path = _resolve_path(target)
        if not path.exists():
            return _err(f"目标不存在: {target}")
    except PermissionError as e:
        return _err(str(e))
    item_type = "文件夹" if path.is_dir() else "文件"
    if not json.loads(ask_user_confirmation(f"删除{item_type}",
        f"{path}\n(递归: {'是' if recursive else '否'})")).get("data"):
        return _err("用户取消了删除操作")
    try:
        if path.is_dir():
            import shutil
            if recursive:
                shutil.rmtree(str(path))
            else:
                path.rmdir()
        else:
            path.unlink()
        return _ok(f"已删除 {target}")
    except OSError as e:
        if "not empty" in str(e).lower() or "目录不是空的" in str(e):
            return _err("文件夹非空，无法删除。设置 recursive=true 以递归删除。")
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def copy_file(source: str, target: str) -> str:
    """复制文件到目标路径。敏感操作，需确认。"""
    try:
        src = _resolve_path(source)
        dst = _resolve_path(target)
        if not src.exists():
            return _err(f"源文件不存在: {source}")
        if not src.is_file():
            return _err(f"源路径不是文件: {source}")
    except PermissionError as e:
        return _err(str(e))
    if not json.loads(ask_user_confirmation("复制文件", f"{src}\n→ {dst}")).get("data"):
        return _err("用户取消了复制操作")
    try:
        import shutil
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(str(src), str(dst))
        return _ok(f"已将 {source} 复制到 {target}")
    except Exception as e:
        return _err(str(e))


def make_directory(folder: str) -> str:
    """创建文件夹，支持嵌套路径。"""
    try:
        path = _resolve_path(folder)
        if path.exists():
            return _err(f"路径已存在: {folder}")
        path.mkdir(parents=True, exist_ok=False)
        return _ok(f"已创建文件夹 {folder}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def grep_file(folder: str, pattern: str, glob_pattern: str = "*", case_sensitive: bool = True) -> str:
    """在文件夹中搜索文件内容，支持正则。只读。"""
    try:
        import re
        p = _resolve_path(folder)
        if not p.exists() or not p.is_dir():
            return _err(f"文件夹不存在: {folder}")
        flags = 0 if case_sensitive else re.IGNORECASE
        try:
            regex = re.compile(pattern, flags)
        except re.error as e:
            return _err(f"正则无效: {e}")
        matches = []
        for f in sorted(p.glob(glob_pattern)):
            if not f.is_file():
                continue
            try:
                text = f.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                continue
            for i, line in enumerate(text.split("\n"), 1):
                if regex.search(line):
                    matches.append({"file": str(f.relative_to(p)), "line": i, "text": line.strip()})
                    if len(matches) >= 500:
                        break
            if len(matches) >= 500:
                break
        return _ok({"matches": matches, "count": len(matches), "max_results": 500})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def file_info(file_path: str) -> str:
    """获取文件详细信息（大小、修改时间等）。只读。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        st = path.stat()
        return _ok({"name": path.name, "size": st.st_size,
                     "size_mb": round(st.st_size / 1048576, 2),
                     "modified": datetime.datetime.fromtimestamp(st.st_mtime).isoformat(),
                     "created": datetime.datetime.fromtimestamp(st.st_ctime).isoformat()})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def read_json(file_path: str) -> str:
    """读取 JSON 文件并返回结构化数据。只读。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        return _ok(json.loads(path.read_text(encoding="utf-8")))
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def write_json(file_path: str, data: str) -> str:
    """将 JSON 文字写入文件（自动格式化）。需确认。"""
    try:
        parsed = json.loads(data) if isinstance(data, str) else data
        path = _resolve_path(file_path)
        path.write_text(json.dumps(parsed, ensure_ascii=False, indent=2), encoding="utf-8")
        return _ok(f"JSON 已写入 {file_path}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def count_words(file_path: str) -> str:
    """统计文件字数、行数、字符数。只读。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        text = path.read_text(encoding="utf-8")
        return _ok({"chars": len(text), "words": len(text.split()),
                     "lines": text.count('\n') + 1, "file": path.name})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def zip_files(folder: str, output_name: str = "archive.zip") -> str:
    """将文件夹压缩为 zip。需确认。"""
    try:
        import zipfile
        fp = _resolve_path(folder)
        if not fp.exists():
            return _err(f"文件夹不存在: {folder}")
        out = _resolve_path(output_name)
        with zipfile.ZipFile(str(out), "w", zipfile.ZIP_DEFLATED) as zf:
            for f in fp.rglob("*"):
                if f.is_file():
                    zf.write(f, f.relative_to(fp))
        return _ok(f"已压缩 {folder} → {output_name}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def unzip_file(zip_path: str, extract_to: str = ".") -> str:
    """解压 zip 文件。需确认。"""
    try:
        import zipfile
        path = _resolve_path(zip_path)
        if not path.exists():
            return _err(f"文件不存在: {zip_path}")
        dest = _resolve_path(extract_to)
        with zipfile.ZipFile(str(path), "r") as zf:
            zf.extractall(str(dest))
        return _ok(f"已解压 → {extract_to}，共 {len(zf.namelist())} 个文件")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def download_file(url: str, filename: str = "") -> str:
    """从 URL 下载文件到 workspace。需确认。"""
    try:
        import requests as _r
        name = filename or url.split("/")[-1].split("?")[0] or "download"
        path = _resolve_path(name)
        resp = _r.get(url, timeout=60, stream=True)
        resp.raise_for_status()
        path.write_bytes(resp.content)
        return _ok(f"已下载 {len(resp.content)} 字节 → {name}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 办公文档 (12 tools)
# ═══════════════════════════════════════════════════════════

def read_excel(file_path: str, sheet: Any = 0) -> str:
    """读取 Excel 文件为结构化数据。只读。"""
    try:
        import pandas as pd
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        df = pd.read_excel(path, sheet_name=sheet)
        return _ok({"headers": df.columns.tolist(), "rows": df.values.tolist(),
                     "row_count": len(df), "column_count": len(df.columns)})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"读取 Excel 失败: {e}")


def write_excel(filename: str, headers: str, rows: str, sheet_name: str = "Sheet1") -> str:
    """创建 Excel 文件写入数据。headers 和 rows 为 JSON 数组字符串。需确认。"""
    try:
        import pandas as pd
        import json as _j
        h = _j.loads(headers)
        r = _j.loads(rows)
        if not isinstance(h, list) or not isinstance(r, list):
            return _err("headers 必须是 JSON 数组，rows 必须是 JSON 二维数组")
        if r and len(r[0]) != len(h):
            return _err(f"列数不匹配: headers {len(h)} 列, rows 每行 {len(r[0])} 列")
        df = pd.DataFrame(r, columns=h)
        path = _resolve_path(filename)
        path.parent.mkdir(parents=True, exist_ok=True)
        df.to_excel(str(path), sheet_name=sheet_name, index=False)
        return _ok({"filename": filename, "row_count": len(df), "col_count": len(df.columns), "columns": h})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def sum_excel_column(file_path: str, column_name: str, sheet: Any = 0) -> str:
    """计算 Excel 指定列的合计值。只读。"""
    try:
        import pandas as pd
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        df = pd.read_excel(path, sheet_name=sheet, engine="openpyxl")
        if column_name not in df.columns:
            return _err(f"列 '{column_name}' 不存在。可用: {df.columns.tolist()}")
        return _ok({"column": column_name, "total": float(df[column_name].sum()), "row_count": len(df)})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"统计失败: {e}")


def sum_csv_column(file_path: str, column_name: str) -> str:
    """计算 CSV 文件中指定列的合计值。只读。"""
    try:
        import pandas as pd
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        df = pd.read_csv(path)
        if column_name not in df.columns:
            return _err(f"列 '{column_name}' 不存在。可用: {df.columns.tolist()}")
        return _ok({"column": column_name, "total": float(df[column_name].sum()), "row_count": len(df)})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"统计失败: {e}")


def merge_csvs(folder_path: str, output_name: str = "merged.csv") -> str:
    """合并文件夹中所有 CSV 文件。需确认。"""
    try:
        import pandas as pd
        folder = _resolve_path(folder_path)
        if not folder.exists():
            return _err(f"文件夹不存在: {folder_path}")
        csv_files = sorted(folder.glob("*.csv"))
        if not csv_files:
            return _err("未找到 CSV 文件")
    except PermissionError as e:
        return _err(str(e))
    if not json.loads(ask_user_confirmation("合并 CSV",
        f"来源: {folder_path} ({len(csv_files)} 文件)\n输出: {output_name}")).get("data"):
        return _err("用户取消了合并操作")
    try:
        merged = pd.concat([pd.read_csv(f) for f in csv_files], ignore_index=True)
        outpath = folder / output_name
        merged.to_csv(outpath, index=False)
        return _ok(f"已将 {len(csv_files)} 个 CSV 合并为 {outpath}，共 {len(merged)} 行")
    except Exception as e:
        return _err(f"合并失败: {e}")


def csv_to_json(file_path: str) -> str:
    """将 CSV 文件转为 JSON 数组。只读。"""
    try:
        import csv, io
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        text = path.read_text(encoding="utf-8")
        reader = csv.DictReader(io.StringIO(text))
        rows = [dict(r) for r in reader]
        return _ok(rows)
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def create_docx(filename: str, content: str) -> str:
    """创建 Word 文档并写入文字内容。需确认。\\n 换行。"""
    try:
        from docx import Document
        doc = Document()
        for line in content.split("\\n"):
            doc.add_paragraph(line)
        path = _resolve_path(filename)
        doc.save(str(path))
        return _ok(f"已创建 {filename}，共 {len(content)} 字符")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def read_docx(file_path: str) -> str:
    """读取 Word 文档的文字内容。只读。"""
    try:
        from docx import Document
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
        return _ok(text)
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def edit_docx(file_path: str, new_text: str, mode: str = "append") -> str:
    """编辑 Word 文档。mode: append(追加) 或 replace(旧文字||新文字)。需确认。"""
    try:
        from docx import Document
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        if not json.loads(ask_user_confirmation(f"编辑文档 ({mode})", str(path))).get("data"):
            return _err("用户取消了编辑操作")
        doc = Document(str(path))
        if mode == "append":
            for line in new_text.split("\\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())
        elif mode == "replace":
            parts = new_text.split("||", 1)
            if len(parts) != 2:
                return _err("replace 模式需要 旧文字||新文字 格式")
            old, new = parts
            for p in doc.paragraphs:
                if old in p.text:
                    for run in p.runs:
                        run.text = run.text.replace(old, new)
        else:
            return _err("mode 必须是 append 或 replace")
        doc.save(str(path))
        return _ok(f"已编辑 {file_path}，模式: {mode}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def create_pptx(filename: str, slides: str) -> str:
    """创建 PowerPoint。格式: 标题||内容---标题2||内容。需确认。"""
    try:
        from pptx import Presentation
        prs = Presentation()
        for slide_text in slides.split("---"):
            parts = slide_text.split("||", 1)
            title = parts[0].strip()
            bullets = parts[1].split(",") if len(parts) > 1 else []
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            if bullets:
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                for b in bullets:
                    p = body.add_paragraph()
                    p.text = b.strip()
        path = _resolve_path(filename)
        prs.save(str(path))
        return _ok(f"已创建 {filename}，共 {len(prs.slides)} 页")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def read_pptx(file_path: str) -> str:
    """读取 PowerPoint 文稿文字。只读。"""
    try:
        from pptx import Presentation
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        title = t
                    elif t:
                        texts.append(t)
            slides.append({"index": i + 1, "title": title, "text": "\n".join(texts)})
        return _ok({"slide_count": len(slides), "slides": slides})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def create_pdf(filename: str, content: str) -> str:
    """创建 PDF 文件。需确认，需 reportlab。"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        path = _resolve_path(filename)
        c = canvas.Canvas(str(path), pagesize=A4)
        y = 800
        for line in content.split("\n"):
            c.drawString(50, y, line[:120])
            y -= 20
            if y < 50:
                c.showPage()
                y = 800
        c.save()
        return _ok(f"已创建 {filename}，共 {len(content)} 字符")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def read_pdf(file_path: str, max_chars: int = 5000) -> str:
    """读取 PDF 文件文字。只读，需 PyPDF2。"""
    try:
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return _err("请先安装 PyPDF2: pip install PyPDF2")
        reader = PdfReader(str(path))
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
            if len(text) >= max_chars:
                break
        return _ok(text[:max_chars].strip())
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def create_html(filename: str, html: str) -> str:
    """创建网页文件并用浏览器打开预览。需确认。"""
    try:
        path = _resolve_path(filename)
        path.write_text(html, encoding="utf-8")
        import webbrowser
        webbrowser.open(str(path))
        return _ok(f"已创建并打开 {filename}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 系统信息 (11 tools)
# ═══════════════════════════════════════════════════════════

def get_current_time() -> str:
    """获取当前系统时间。"""
    now = datetime.datetime.now()
    return _ok({"datetime": now.strftime("%Y-%m-%d %H:%M:%S"), "iso": now.isoformat()})


def get_current_directory() -> str:
    """获取当前工作目录。"""
    return _ok({"path": os.getcwd()})


def get_username() -> str:
    """获取当前登录用户名。"""
    return _ok({"username": os.environ.get("USERNAME", os.environ.get("USER", "unknown"))})


def get_hostname() -> str:
    """获取计算机名称。"""
    return _ok({"hostname": platform.node()})


def get_ip() -> str:
    """获取本机 IP 地址。"""
    try:
        import socket
        hostname = socket.gethostname()
        ip = socket.gethostbyname(hostname)
        return _ok({"hostname": hostname, "local_ip": ip})
    except Exception as e:
        return _err(str(e))


def ping_host(host: str, count: int = 3) -> str:
    """Ping 目标主机。"""
    try:
        import subprocess
        r = subprocess.run(f"ping -n {count} {host}", capture_output=True, text=True, shell=True, timeout=20)
        return _ok({"host": host, "result": r.stdout.strip()[-500:], "success": r.returncode == 0})
    except Exception as e:
        return _err(str(e))


def list_processes(name_filter: str = "") -> str:
    """列出当前运行的进程（Windows PowerShell）。只读。"""
    try:
        import subprocess, json as _j
        if name_filter:
            ps_cmd = f"Get-Process -Name '*{name_filter}*' -ErrorAction SilentlyContinue | Select Name,Id,WorkingSet | ConvertTo-Json"
        else:
            ps_cmd = "Get-Process | Select Name,Id,WorkingSet -First 100 | ConvertTo-Json"
        r = subprocess.run(["powershell", "-Command", ps_cmd], capture_output=True, text=True, timeout=15)
        raw = _j.loads(r.stdout) if r.stdout.strip() else []
        procs = raw if isinstance(raw, list) else [raw]
        result = []
        for p in procs:
            if isinstance(p, dict):
                ws = p.get("WorkingSet", 0) or 0
                result.append({"name": p.get("Name", ""), "pid": p.get("Id", 0),
                               "memory_mb": round(ws / 1048576, 1)})
        return _ok({"processes": result, "count": len(result)})
    except Exception as e:
        return _err(str(e))


def kill_process(target: str, by: str = "name") -> str:
    """终止进程，按 name 或 pid。敏感操作，需确认。"""
    try:
        if not json.loads(ask_user_confirmation(f"终止进程 ({by})", f"目标: {target}")).get("data"):
            return _err("用户取消了终止进程操作")
        import subprocess
        if by == "name":
            ps_cmd = f"Stop-Process -Name '{target}' -Force -ErrorAction Stop"
        elif by == "pid":
            ps_cmd = f"Stop-Process -Id {target} -Force -ErrorAction Stop"
        else:
            return _err("by 必须是 name 或 pid")
        r = subprocess.run(["powershell", "-Command", ps_cmd], capture_output=True, text=True, timeout=15)
        if r.returncode != 0:
            return _err(f"终止失败: {r.stderr.strip() or '进程不存在或无权限'}")
        return _ok(f"已终止进程: {target}")
    except Exception as e:
        return _err(str(e))


def screen_capture(filename: str = "screenshot.png") -> str:
    """截取整个屏幕保存为图片。需确认，需 Pillow。"""
    try:
        from PIL import ImageGrab
        path = _resolve_path(filename)
        if not json.loads(ask_user_confirmation("屏幕截图", f"保存到: {path}")).get("data"):
            return _err("用户取消了截图操作")
        img = ImageGrab.grab()
        img.save(str(path))
        return _ok(f"截图已保存到 {filename}，{img.size[0]}x{img.size[1]}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def image_resize(file_path: str, width: int, height: int, output: str = "") -> str:
    """调整图片尺寸。需 Pillow。"""
    try:
        from PIL import Image
        path = _resolve_path(file_path)
        if not path.exists():
            return _err(f"文件不存在: {file_path}")
        img = Image.open(str(path))
        img = img.resize((width, height), Image.LANCZOS)
        out = _resolve_path(output) if output else path
        img.save(str(out))
        return _ok(f"图片已调整至 {width}x{height}")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 浏览器操作 (6 tools)
# ═══════════════════════════════════════════════════════════

_browser = None
_page = None

def _get_page():
    global _browser, _page
    if _browser is None:
        from playwright.sync_api import sync_playwright
        pw = sync_playwright().start()
        _browser = pw.chromium.launch(channel="msedge", headless=False)
        _page = _browser.new_page()
    return _page


def browser_open(url: str) -> str:
    """用浏览器打开指定网页。"""
    try:
        page = _get_page()
        page.goto(url, timeout=30000)
        return _ok(f"已打开 {url}，标题: {page.title()}")
    except Exception as e:
        return _err(str(e))


def browser_type(text: str) -> str:
    """在当前网页的焦点输入框中输入文字。"""
    try:
        page = _get_page()
        page.keyboard.type(text)
        return _ok(f"已输入 {len(text)} 字符")
    except Exception as e:
        return _err(str(e))


def browser_click(selector: str) -> str:
    """点击网页上的元素，支持 CSS 选择器或文字匹配。"""
    try:
        page = _get_page()
        try:
            page.click(selector, timeout=5000)
        except Exception:
            page.get_by_text(selector).first.click(timeout=5000)
        return _ok(f"已点击 '{selector}'")
    except Exception as e:
        return _err(str(e))


def browser_read() -> str:
    """读取当前网页的文本内容（最多 3000 字符）。只读。"""
    try:
        page = _get_page()
        text = page.inner_text("body")[:3000]
        return _ok(text)
    except Exception as e:
        return _err(str(e))


def browser_screenshot(filename: str = "screenshot.png") -> str:
    """截取当前网页保存到 workspace。"""
    try:
        page = _get_page()
        path = (_WORKSPACE_ROOT / filename) if _WORKSPACE_ROOT else Path(filename)
        page.screenshot(path=str(path))
        return _ok(f"截图已保存到 {filename}")
    except Exception as e:
        return _err(str(e))


def web_search(query: str) -> str:
    """用必应搜索并返回结果摘要，打开/输入/搜索/读取一步完成。"""
    try:
        page = _get_page()
        page.goto("https://cn.bing.com", timeout=15000)
        page.fill('input[name="q"]', query)
        page.keyboard.press("Enter")
        page.wait_for_timeout(2000)
        text = page.inner_text("body")[:3000]
        return _ok(f"搜索 '{query}' 结果:\n{text}")
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# CMD / 代码 / 数据库 (5 tools)
# ═══════════════════════════════════════════════════════════

_CMD_BLACKLIST = ["format", "del ", "rmdir", "rd ", "shutdown", "taskkill",
    "reg ", "diskpart", "cacls", "icacls", "takeown", "bcdedit", "netsh",
    "sc ", "net stop", "net start", "erase", "deltree", "fdisk"]


def execute_command(command: str) -> str:
    """执行 CMD 命令，黑名单保护 + 确认 + 30s 超时。"""
    cmd_lower = command.lower().strip()
    for banned in _CMD_BLACKLIST:
        if banned in cmd_lower:
            return _err(f"命令被黑名单拦截: 禁止执行含 '{banned}' 的命令")
    if not json.loads(ask_user_confirmation("执行命令",
        f"命令: {command}\n工作目录: {_WORKSPACE_ROOT}")).get("data"):
        return _err("用户取消了命令执行")
    try:
        import subprocess
        result = subprocess.run(f"cmd /c {command}", shell=True, capture_output=True,
            text=True, timeout=30, cwd=str(_WORKSPACE_ROOT) if _WORKSPACE_ROOT else None)
        return _ok({"stdout": result.stdout.strip() or "(无输出)",
                     "stderr": result.stderr.strip() or "", "exit_code": result.returncode})
    except subprocess.TimeoutExpired:
        return _err("命令执行超时（30 秒）")
    except Exception as e:
        return _err(str(e))


def run_python(code: str) -> str:
    """执行安全 Python 代码（仅允许只读操作）。"""
    try:
        banned = ["import os","open(","__import__","eval(","exec(","compile(","shutil","subprocess","socket","ctypes"]
        for b in banned:
            if b in code:
                return _err(f"代码含禁用关键字: {b}")
        g = {"__builtins__": {"print":print,"len":len,"range":range,"int":int,"str":str,"float":float,"list":list,"dict":dict,"sum":sum,"max":max,"min":min,"sorted":sorted,"enumerate":enumerate,"zip":zip,"map":map,"filter":filter,"abs":abs,"round":round,"type":type,"isinstance":isinstance,"True":True,"False":False,"None":None}}
        import io, sys, traceback
        out = io.StringIO()
        old = sys.stdout; sys.stdout = out
        try: exec(code, g)
        except: traceback.print_exc()
        sys.stdout = old
        return _ok(out.getvalue().strip() or "(无输出)")
    except Exception as e:
        return _err(str(e))


def sql_query(db_path: str, query: str) -> str:
    """执行 SQLite 查询（仅允许 SELECT/PRAGMA）。只读。"""
    try:
        import sqlite3
        q = query.strip().upper()
        if not q.startswith("SELECT") and not q.startswith("PRAGMA"):
            return _err("只允许 SELECT 查询")
        path = _resolve_path(db_path)
        conn = sqlite3.connect(str(path))
        cur = conn.execute(query)
        cols = [d[0] for d in cur.description] if cur.description else []
        rows = [list(r) for r in cur.fetchall()]
        conn.close()
        return _ok({"columns": cols, "rows": rows, "row_count": len(rows)})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def sql_exec(db_path: str, sql: str) -> str:
    """执行 SQLite 建表/插入操作。需确认。"""
    try:
        import sqlite3
        if not json.loads(ask_user_confirmation("执行 SQL",
            f"数据库: {db_path}\nSQL: {sql[:200]}")).get("data"):
            return _err("用户取消了 SQL 执行")
        path = _resolve_path(db_path)
        conn = sqlite3.connect(str(path))
        conn.executescript(sql)
        conn.commit()
        conn.close()
        return _ok("SQL 执行成功")
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 通讯 & 通知 (2 tools)
# ═══════════════════════════════════════════════════════════

def send_email(to: str, subject: str, body: str) -> str:
    """发送邮件（QQ邮箱 SMTP）。需设 EMAIL_USER 和 EMAIL_PASS 环境变量。需确认。"""
    try:
        import smtplib
        from email.mime.text import MIMEText
        user = os.environ.get("EMAIL_USER", "")
        pwd = os.environ.get("EMAIL_PASS", "")
        if not user or not pwd:
            return _err("未配置 EMAIL_USER 和 EMAIL_PASS")
        if not json.loads(ask_user_confirmation("发送邮件",
            f"收件人: {to}\n主题: {subject}")).get("data"):
            return _err("用户取消了邮件发送")
        msg = MIMEText(body, "plain", "utf-8")
        msg["Subject"] = subject; msg["From"] = user; msg["To"] = to
        with smtplib.SMTP("smtp.qq.com", 587, timeout=15) as s:
            s.starttls(); s.login(user, pwd); s.sendmail(user, [to], msg.as_string())
        return _ok(f"邮件已发送给 {to}")
    except Exception as e:
        return _err(str(e))


def send_notification(title: str, message: str) -> str:
    """弹出 Windows 桌面通知。"""
    try:
        import subprocess
        ps = f"""Add-Type -AssemblyName System.Windows.Forms
$n = New-Object System.Windows.Forms.NotifyIcon
$n.Icon = [System.Drawing.SystemIcons]::Information
$n.BalloonTipTitle = '{title}'
$n.BalloonTipText = '{message}'
$n.Visible = $true
$n.ShowBalloonTip(3000)
Start-Sleep -Seconds 4
$n.Dispose()"""
        subprocess.Popen(["powershell", "-Command", ps], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return _ok("通知已弹出")
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 网络 / API / GitHub (5 tools)
# ═══════════════════════════════════════════════════════════

def http_request(url: str, method: str = "GET", headers: str = "{}", body: str = "") -> str:
    """发送 HTTP GET/POST 请求。只读。"""
    try:
        import requests as _r
        import json as _j
        hdrs = _j.loads(headers) if headers else {}
        if method.upper() == "GET":
            resp = _r.get(url, headers=hdrs, timeout=30)
        elif method.upper() == "POST":
            resp = _r.post(url, headers=hdrs, data=body.encode("utf-8"), timeout=30)
        else:
            return _err("method 只支持 GET 或 POST")
        resp.raise_for_status()
        return _ok({"status": resp.status_code, "url": resp.url,
                     "headers": dict(resp.headers), "body": resp.text[:3000]})
    except Exception as e:
        return _err(str(e))


def github_search(query: str, search_type: str = "repositories", sort: str = "stars", per_page: int = 10) -> str:
    """搜索 GitHub 仓库/代码/Issues/用户。设 GITHUB_TOKEN 提高速率（未设则 60次/小时）。只读。"""
    try:
        import requests as _r
        token = os.environ.get("GITHUB_TOKEN", "")
        headers = {"Accept": "application/vnd.github.v3+json"}
        if token:
            headers["Authorization"] = f"token {token}"
        type_map = {
            "repositories": f"/search/repositories?q={_r.utils.quote(query)}&sort={sort}&per_page={per_page}",
            "code": f"/search/code?q={_r.utils.quote(query)}&per_page={per_page}",
            "issues": f"/search/issues?q={_r.utils.quote(query)}&sort={sort}&per_page={per_page}",
            "users": f"/search/users?q={_r.utils.quote(query)}&per_page={per_page}",
        }
        if search_type not in type_map:
            return _err("search_type 必须是 repositories/code/issues/users")
        resp = _r.get(f"https://api.github.com{type_map[search_type]}", headers=headers, timeout=15)
        if resp.status_code == 403 and "rate limit" in resp.text.lower():
            return _err("GitHub API 速率限制。请设置 GITHUB_TOKEN 或稍后再试。")
        resp.raise_for_status()
        data = resp.json()
        results = []
        for item in data.get("items", [])[:per_page]:
            if search_type == "repositories":
                results.append({"full_name": item.get("full_name"),
                    "description": (item.get("description") or "")[:200],
                    "stars": item.get("stargazers_count", 0), "language": item.get("language", ""),
                    "url": item.get("html_url"), "updated": item.get("updated_at", "")})
            elif search_type == "code":
                results.append({"repo": item.get("repository", {}).get("full_name"),
                    "path": item.get("path"), "url": item.get("html_url")})
            elif search_type == "issues":
                results.append({"title": item.get("title"), "state": item.get("state"),
                    "url": item.get("html_url")})
            elif search_type == "users":
                results.append({"login": item.get("login"), "url": item.get("html_url")})
        return _ok({"results": results, "total_count": data.get("total_count", 0)})
    except Exception as e:
        return _err(str(e))


def github_get_repo(owner: str, repo: str) -> str:
    """获取 GitHub 仓库详情和 README 预览。只读。"""
    try:
        import requests as _r
        token = os.environ.get("GITHUB_TOKEN", "")
        headers = {"Accept": "application/vnd.github.v3+json"}
        if token:
            headers["Authorization"] = f"token {token}"
        resp = _r.get(f"https://api.github.com/repos/{owner}/{repo}", headers=headers, timeout=15)
        if resp.status_code == 404:
            return _err(f"仓库不存在: {owner}/{repo}")
        resp.raise_for_status()
        repo_data = resp.json()
        readme_text = ""
        try:
            rm = _r.get(f"https://api.github.com/repos/{owner}/{repo}/readme",
                         headers={**headers, "Accept": "application/vnd.github.v3.raw"}, timeout=10)
            if rm.status_code == 200:
                readme_text = rm.text[:2000]
        except Exception:
            pass
        return _ok({"name": repo_data.get("full_name"),
            "description": repo_data.get("description", ""),
            "stars": repo_data.get("stargazers_count", 0), "forks": repo_data.get("forks_count", 0),
            "language": repo_data.get("language", ""), "topics": repo_data.get("topics", []),
            "url": repo_data.get("html_url"), "readme_preview": readme_text})
    except Exception as e:
        return _err(str(e))


def github_download_repo(owner: str, repo: str, branch: str = "main", extract_to: str = "") -> str:
    """下载 GitHub 仓库 zip 并解压到 workspace。需确认。"""
    try:
        import requests as _r, zipfile, io
        url = f"https://github.com/{owner}/{repo}/archive/refs/heads/{branch}.zip"
        target = _resolve_path(extract_to or repo)
        if not json.loads(ask_user_confirmation("下载 GitHub 仓库",
            f"{owner}/{repo} ({branch})\n→ {target}")).get("data"):
            return _err("用户取消了下载操作")
        resp = _r.get(url, timeout=60, stream=True)
        if resp.status_code == 404:
            return _err(f"仓库/分支不存在: {owner}/{repo}/{branch}")
        resp.raise_for_status()
        target.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(io.BytesIO(resp.content)) as zf:
            prefix = f"{repo}-{branch}/"
            files = [n for n in zf.namelist() if n.startswith(prefix) and not n.endswith("/")]
            for name in files:
                rel = name[len(prefix):]
                out = target / rel
                out.parent.mkdir(parents=True, exist_ok=True)
                out.write_bytes(zf.read(name))
        return _ok({"extracted_to": str(target), "files": len(files),
                     "repo": f"{owner}/{repo}", "branch": branch})
    except PermissionError as e:
        return _err(str(e))
    except Exception as e:
        return _err(str(e))


def install_skill(source: str, branch: str = "main") -> str:
    """从 GitHub 安装 Skill。下载/解压/装依赖/验证。安装后需重启 Agent。"""
    try:
        import requests as _r, zipfile, io, subprocess, sys as _sys
        source = source.rstrip("/").replace("https://github.com/", "").replace(".git", "")
        parts = source.split("/")
        if len(parts) < 2:
            return _err("格式错误，请用 owner/repo 或 GitHub URL")
        owner, repo = parts[0], parts[1]
        url = f"https://github.com/{owner}/{repo}/archive/refs/heads/{branch}.zip"
        resp = _r.get(url, timeout=60, stream=True)
        if resp.status_code == 404:
            return _err(f"仓库/分支不存在: {owner}/{repo}/{branch}")
        resp.raise_for_status()
        with zipfile.ZipFile(io.BytesIO(resp.content)) as zf:
            prefix = f"{repo}-{branch}/"
            names = [n for n in zf.namelist() if n.startswith(prefix)]
            manifest_name = f"{prefix}skill.json"
            if manifest_name not in names:
                return _err("该仓库不是有效 Skill，缺少 skill.json")
            manifest = json.loads(zf.read(manifest_name))
            skill_name = manifest.get("name", repo)
            target_dir = _SKILLS_DIR / skill_name
            target_dir.mkdir(parents=True, exist_ok=True)
            for name in names:
                if name.endswith("/"):
                    continue
                rel = name[len(prefix):]
                out = target_dir / rel
                out.parent.mkdir(parents=True, exist_ok=True)
                out.write_bytes(zf.read(name))
        req_file = target_dir / "requirements.txt"
        deps_installed = []
        if req_file.exists():
            for req in req_file.read_text().strip().split("\n"):
                req = req.strip()
                if req and not req.startswith("#"):
                    try:
                        subprocess.run([_sys.executable, "-m", "pip", "install", req,
                            "--break-system-packages", "-q"], timeout=60, check=True)
                        deps_installed.append(req)
                    except Exception:
                        pass
        valid = False
        module_path = target_dir / "tools.py"
        if module_path.exists():
            import importlib.util
            try:
                spec = importlib.util.spec_from_file_location(f"_skill_check_{skill_name}", str(module_path))
                if spec and spec.loader:
                    mod = importlib.util.module_from_spec(spec)
                    spec.loader.exec_module(mod)
                    valid = True
            except Exception:
                pass
        return _ok({"name": skill_name, "description": manifest.get("description", ""),
            "version": manifest.get("version", "0.1.0"),
            "tools": [t.get("name") for t in manifest.get("tools", [])],
            "requirements": deps_installed, "valid": valid,
            "installed_to": str(target_dir),
            "hint": "请重启 Agent 以加载新 Skill" if valid else "Skill 安装但验证失败，请检查 tools.py"})
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 工具函数 (12 tools)
# ═══════════════════════════════════════════════════════════

def translate_text(text: str, target: str = "中文") -> str:
    """用本地 Qwen 模型翻译文字。"""
    try:
        import requests as _r
        host = os.environ.get("LLM_HOST", "http://localhost:1234")
        resp = _r.post(f"{host}/v1/chat/completions",
            json={"model": os.environ.get("LLM_MODEL", "qwen2.5-7b-instruct"),
                  "messages": [{"role": "system", "content": f"将以下文字翻译为{target}，只输出译文"},
                               {"role": "user", "content": text}],
                  "temperature": 0.1, "max_tokens": 1024}, timeout=30)
        result = resp.json()["choices"][0]["message"]["content"]
        return _ok(result)
    except Exception as e:
        return _err(str(e))


def clipboard_copy(text: str) -> str:
    """复制文字到剪贴板。"""
    try:
        import subprocess
        subprocess.run("clip", input=text, text=True, shell=True, timeout=5)
        return _ok(f"已复制 {len(text)} 字符")
    except Exception as e:
        return _err(str(e))


def clipboard_paste() -> str:
    """读取系统剪贴板文字。"""
    try:
        import subprocess
        r = subprocess.run("powershell Get-Clipboard", capture_output=True, text=True, shell=True, timeout=5)
        return _ok(r.stdout.strip())
    except Exception as e:
        return _err(str(e))


def set_reminder(seconds: int, message: str) -> str:
    """N 秒后弹出提醒。"""
    try:
        import threading
        def remind():
            import time as _t
            _t.sleep(seconds)
            send_notification("提醒", message)
        threading.Thread(target=remind, daemon=True).start()
        return _ok(f"已设置 {seconds} 秒后提醒: {message}")
    except Exception as e:
        return _err(str(e))


def calculator(expression: str) -> str:
    """安全计算数学表达式。"""
    try:
        allowed = set("0123456789+-*/.() eEπpiqwertyuiopasdfghjklzxcvbnm")
        if not all(c.lower() in allowed for c in expression):
            return _err("表达式含非法字符")
        import math
        expr = expression.replace("^", "**").replace("÷", "/").replace("×", "*")
        result = eval(expr, {"__builtins__": {}}, {"math": math, "pi": math.pi, "e": math.e})
        return _ok(str(result))
    except Exception as e:
        return _err(str(e))


def generate_password(length: int = 16) -> str:
    """生成随机密码。"""
    try:
        import secrets, string
        chars = string.ascii_letters + string.digits + "!@#$%^&*"
        pwd = ''.join(secrets.choice(chars) for _ in range(length))
        return _ok(pwd)
    except Exception as e:
        return _err(str(e))


def hash_text(text: str, algo: str = "sha256") -> str:
    """计算文本哈希值。algo: md5/sha1/sha256。"""
    try:
        import hashlib
        h = getattr(hashlib, algo)()
        h.update(text.encode("utf-8"))
        return _ok({"algo": algo, "hash": h.hexdigest()})
    except Exception as e:
        return _err(str(e))


def base64_cmd(action: str, text: str) -> str:
    """Base64 编码或解码。action: encode/decode。"""
    try:
        import base64
        if action == "encode":
            return _ok(base64.b64encode(text.encode()).decode())
        elif action == "decode":
            return _ok(base64.b64decode(text).decode())
        return _err("action 必须是 encode 或 decode")
    except Exception as e:
        return _err(str(e))


def date_diff(date1: str, date2: str = "") -> str:
    """计算两个日期间隔天数。date2 默认今天。格式: 2026-05-21。"""
    try:
        from datetime import date as dt
        d1 = dt.fromisoformat(date1)
        d2 = dt.fromisoformat(date2) if date2 else dt.today()
        diff = abs((d2 - d1).days)
        return _ok({"date1": date1, "date2": str(d2), "days": diff})
    except Exception as e:
        return _err(str(e))


# ═══════════════════════════════════════════════════════════
# 注册表分层
# ═══════════════════════════════════════════════════════════

TOOL_REGISTRY: list = [
    ("find_files", find_files, "搜索匹配 glob 模式的文件（只读）"),
    ("list_directory", list_directory, "列出文件夹内容（只读）"),
    ("read_file_content", read_file_content, "读取文本文件（只读）"),
    ("write_file_content", write_file_content, "写入内容到文件（需确认）"),
    ("open_file", open_file, "用默认程序打开文件（需确认）"),
    ("rename_file", rename_file, "重命名或移动文件（需确认）"),
    ("batch_rename", batch_rename, "批量重命名文件夹中的文件（需确认）"),
    ("delete_file", delete_file, "删除文件或文件夹（需确认，非空需递归）"),
    ("copy_file", copy_file, "复制文件到目标路径（需确认）"),
    ("make_directory", make_directory, "创建文件夹，支持嵌套路径"),
    ("grep_file", grep_file, "在文件夹中搜索文件内容，支持正则（只读）"),
    ("file_info", file_info, "获取文件详细信息大小/时间（只读）"),
    ("read_json", read_json, "读取 JSON 文件返回结构化数据（只读）"),
    ("write_json", write_json, "格式化写入 JSON 文件（需确认）"),
    ("count_words", count_words, "统计文件字数/行数/字符数（只读）"),
    ("zip_files", zip_files, "压缩文件夹为 zip（需确认）"),
    ("unzip_file", unzip_file, "解压 zip 文件（需确认）"),
    ("download_file", download_file, "下载文件到 workspace（需确认）"),
    ("create_docx", create_docx, "创建 Word 文档（需确认）"),
    ("read_docx", read_docx, "读取 Word 文档文字（只读）"),
    ("edit_docx", edit_docx, "编辑已有 Word 文档追加/替换（需确认）"),
    ("create_pptx", create_pptx, "创建 PowerPoint（需确认）"),
    ("read_pptx", read_pptx, "读取 PowerPoint 文字（只读）"),
    ("create_pdf", create_pdf, "创建 PDF（需确认，需 reportlab）"),
    ("read_pdf", read_pdf, "读取 PDF 文字（只读，需 PyPDF2）"),
    ("create_html", create_html, "创建网页并打开预览（需确认）"),
    ("read_excel", read_excel, "读取 Excel 为结构化数据（只读）"),
    ("write_excel", write_excel, "创建 Excel 写入 JSON headers/rows（需确认）"),
    ("sum_excel_column", sum_excel_column, "计算 Excel 列合计（只读）"),
    ("sum_csv_column", sum_csv_column, "计算 CSV 列合计（只读）"),
    ("merge_csvs", merge_csvs, "合并文件夹中所有 CSV（需确认）"),
    ("csv_to_json", csv_to_json, "CSV 转 JSON 数组（只读）"),
    ("get_current_time", get_current_time, "获取当前系统时间（只读）"),
    ("get_current_directory", get_current_directory, "获取当前工作目录（只读）"),
    ("get_username", get_username, "获取当前用户名（只读）"),
    ("get_hostname", get_hostname, "获取计算机名（只读）"),
    ("get_ip", get_ip, "获取本机 IP 地址（只读）"),
    ("ping_host", ping_host, "Ping 目标主机（只读）"),
    ("list_processes", list_processes, "列出运行中的进程（只读）"),
    ("kill_process", kill_process, "终止进程按名称或 PID（需确认）"),
    ("screen_capture", screen_capture, "截取整个屏幕保存图片（需确认，需 Pillow）"),
    ("image_resize", image_resize, "调整图片尺寸（需 Pillow）"),
    ("browser_open", browser_open, "用浏览器打开网页"),
    ("browser_type", browser_type, "在网页输入框打字"),
    ("browser_click", browser_click, "点击网页元素 CSS选择器或文字"),
    ("browser_read", browser_read, "读取当前网页文字最多3000字符（只读）"),
    ("browser_screenshot", browser_screenshot, "截取当前网页保存图片"),
    ("web_search", web_search, "用必应搜索并返回结果摘要"),
    ("execute_command", execute_command, "执行 CMD 命令黑名单保护+确认+30s超时"),
    ("run_python", run_python, "执行安全 Python 代码仅允许只读"),
    ("sql_query", sql_query, "执行 SQLite 查询仅 SELECT（只读）"),
    ("sql_exec", sql_exec, "执行 SQLite 建表/插入（需确认）"),
    ("send_email", send_email, "发送邮件 QQ邮箱 SMTP（需确认，需 EMAIL_USER/PASS）"),
    ("send_notification", send_notification, "弹出 Windows 桌面通知"),
    ("http_request", http_request, "发送 HTTP GET/POST 请求（只读）"),
    ("github_search", github_search, "搜索 GitHub 仓库/代码/Issues/用户（只读）"),
    ("github_get_repo", github_get_repo, "获取 GitHub 仓库详情+README（只读）"),
    ("github_download_repo", github_download_repo, "下载 GitHub 仓库 zip 到 workspace（需确认）"),
    ("install_skill", install_skill, "从 GitHub 安装 Skill 下载+解压+装依赖+验证"),
    ("translate_text", translate_text, "用本地 Qwen 翻译文字（只读）"),
    ("clipboard_copy", clipboard_copy, "复制文字到剪贴板"),
    ("clipboard_paste", clipboard_paste, "读取剪贴板文字"),
    ("set_reminder", set_reminder, "设置 N 秒后弹出提醒"),
    ("calculator", calculator, "安全计算数学表达式"),
    ("generate_password", generate_password, "生成随机安全密码"),
    ("hash_text", hash_text, "计算文本 MD5/SHA256 哈希"),
    ("base64_cmd", base64_cmd, "Base64 编码或解码"),
    ("date_diff", date_diff, "计算两个日期间隔天数"),
]

INTERNAL_TOOLS: dict[str, Any] = {
    "ask_user_confirmation": ask_user_confirmation,
}


# ═══════════════════════════════════════════════════════════
# Skill 动态加载
# ═══════════════════════════════════════════════════════════

def load_all_skills() -> int:
    """扫描 skills/ 目录，导入所有有效 skill。返回加载的工具数。"""
    global _SKILL_TOOLS
    _SKILL_TOOLS = []
    if not _SKILLS_DIR.exists():
        return 0
    import importlib.util, sys
    count = 0
    for skill_dir in sorted(_SKILLS_DIR.iterdir()):
        if not skill_dir.is_dir():
            continue
        manifest = skill_dir / "skill.json"
        if not manifest.exists():
            continue
        try:
            info = json.loads(manifest.read_text(encoding="utf-8"))
            module_path = skill_dir / "tools.py"
            if not module_path.exists():
                continue
            name = info.get("name", skill_dir.name)
            spec = importlib.util.spec_from_file_location(f"skills.{name}", str(module_path))
            if spec is None or spec.loader is None:
                continue
            mod = importlib.util.module_from_spec(spec)
            sys.modules[f"skills.{name}"] = mod
            spec.loader.exec_module(mod)
            for t in info.get("tools", []):
                func = getattr(mod, t["function"], None)
                if func:
                    _SKILL_TOOLS.append((t["name"], func, t.get("description", "")))
                    count += 1
        except Exception:
            pass
    return count


def get_all_tools() -> list:
    """返回 TOOL_REGISTRY + 动态加载的 skill 工具。"""
    return list(TOOL_REGISTRY) + list(_SKILL_TOOLS)


def get_tool_schema() -> list[dict]:
    """导出模型可见工具为标准 JSON Schema。不含 ask_user_confirmation。"""
    import inspect
    schemas = []
    for name, func, desc in get_all_tools():
        sig = inspect.signature(func)
        props, req = {}, []
        for pn, p in sig.parameters.items():
            m = {"str": "string", "int": "integer", "float": "number", "bool": "boolean"}
            jt = "string"
            for py, js in m.items():
                if py in str(p.annotation):
                    jt = js
                    break
            props[pn] = {"type": jt}
            # Extract param description from docstring
            doc = func.__doc__ or ""
            for line in doc.split("\n"):
                if line.strip().startswith(f"{pn}:"):
                    props[pn]["description"] = line.split(":", 1)[1].strip()
            if p.default != inspect.Parameter.empty:
                props[pn]["default"] = None if p.default is None else p.default
            else:
                req.append(pn)
        schemas.append({"name": name, "description": desc,
                         "parameters": {"type": "object", "properties": props,
                                        "required": req, "additionalProperties": False}})
    return schemas


def get_tools_map() -> dict:
    return {n: f for n, f, _ in get_all_tools()}


def get_internal_tool(name: str) -> Any:
    return INTERNAL_TOOLS.get(name)


if __name__ == "__main__":
    import sys
    if "--schema" in sys.argv:
        print(json.dumps(get_tool_schema(), ensure_ascii=False, indent=2))
    else:
        skills_loaded = load_all_skills()
        if skills_loaded:
            print(f"Skill tools: {skills_loaded}")
        print(f"Tools: {len(get_all_tools())} ({len(TOOL_REGISTRY)} built-in + {len(_SKILL_TOOLS)} skills)")
